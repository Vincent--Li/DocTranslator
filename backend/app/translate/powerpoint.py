import threading
import pptx
from . import to_translate
from . import common
import os
import sys
import time
import datetime

def start(trans):
    # 允许的最大线程
    threads=trans['threads']
    if threads is None or int(threads)<0:
        max_threads=10
    else:
        max_threads=int(threads)
    # 当前执行的索引位置
    run_index=0
    start_time = datetime.datetime.now()
    wb = pptx.Presentation(trans['file_path']) 
    print(trans['file_path'])
    slides = wb.slides

    # 提取文本
    texts=[]
    for slide in slides:
        for shape in slide.shapes:
            texts.extend(pptx_extract_text_from_shape(shape))
    max_run=max_threads if len(texts)>max_threads else len(texts)
    before_active_count=threading.active_count()
    event=threading.Event()
    while run_index<=len(texts)-1:
        if threading.active_count()<max_run+before_active_count:
            if not event.is_set():
                thread = threading.Thread(target=to_translate.get,args=(trans,event,texts,run_index))
                thread.start()
                run_index+=1
            else:
                return False
    
    while True:
        complete=True
        for text in texts:
            if not text['complete']:
                complete=False
        if complete:
            break
        else:
            time.sleep(1)

    # 回写翻译
    text_count=0
    for slide in slides:
        for shape in slide.shapes:
            text_count += pptx_write_translation_to_shape(shape, texts)

    wb.save(trans['target_file'])
    end_time = datetime.datetime.now()
    spend_time=common.display_spend(start_time, end_time)
    to_translate.complete(trans,text_count,spend_time)
    return True


def pptx_extract_text_from_shape(shape):
    ### pptx文件特定方法，提取pptx中待翻译的文本块
    texts = []
    if shape.shape_type == 6:  # GroupShape 类型
        for s in shape.shapes:
            texts.extend(pptx_extract_text_from_shape(s))
    elif shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                # 纯符号不用翻译
                if run.text.strip():
                    texts.append({"text":run.text, "complete":False})
    return texts

def pptx_write_translation_to_shape6(shape, texts):
    ### pptx文件特定方法，提取pptx中待翻译的文本块
    tmp_text_count = 0
    if shape.shape_type == 6:  # GroupShape 类型
        for s in shape.shapes:
            tmp_text_count += pptx_write_translation_to_shape6(s, texts)
    elif shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                # 纯符号不用翻译
                if run.text.strip():
                    item = texts.pop(0)
                    run.text = item['text']
                    tmp_text_count+=item['count']
    return tmp_text_count
        
def pptx_write_translation_to_shape(shape, texts):
    ### 回写翻译结果到pptx文件
    tmp_text_count = 0
    if shape.shape_type == 6:  # GroupShape 类型
        for s in shape.shapes:
            tmp_text_count += pptx_write_translation_to_shape6(s, texts)
    elif shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text and run.text.strip():
                    item=texts.pop(0)
                    run.text=item['text']
                    tmp_text_count+=item['count']
    return tmp_text_count
